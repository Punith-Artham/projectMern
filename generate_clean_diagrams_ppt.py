from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


FILL = RGBColor(239, 246, 255)
LINE = RGBColor(59, 130, 246)
ACCENT = RGBColor(79, 70, 229)


def add_box(slide, x, y, w, h, text, size=13):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = FILL
    s.line.color.rgb = LINE
    tf = s.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    return s


def add_arrow(slide, x, y, w, h, text=""):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid()
    a.fill.fore_color.rgb = ACCENT
    a.line.color.rgb = ACCENT
    if text:
        t = a.text_frame
        t.text = text
        t.paragraphs[0].font.size = Pt(10)
        t.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        t.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_down(slide, x, y, w, h):
    d = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    d.fill.solid()
    d.fill.fore_color.rgb = ACCENT
    d.line.color.rgb = ACCENT


def add_title(prs, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = title
    s.placeholders[1].text = subtitle


def architecture(prs):
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Architecture Diagram"
    add_box(s, 0.6, 1.6, 2.0, 0.8, "User Browser")
    add_box(s, 3.0, 1.6, 2.3, 0.8, "React Frontend")
    add_box(s, 5.8, 1.6, 2.4, 0.8, "Express API")
    add_box(s, 8.7, 1.6, 1.7, 0.8, "MongoDB")
    add_box(s, 5.8, 3.4, 2.4, 0.8, "Flask AI Service")
    add_box(s, 5.8, 5.0, 2.4, 0.8, "Vision + ML Engine")
    add_arrow(s, 2.65, 1.9, 0.3, 0.25)
    add_arrow(s, 5.35, 1.9, 0.35, 0.25)
    add_arrow(s, 8.3, 1.9, 0.35, 0.25)
    add_down(s, 6.8, 2.5, 0.45, 0.8)
    add_down(s, 6.8, 4.3, 0.45, 0.6)


def use_case(prs):
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Use Case Diagram"
    add_box(s, 0.5, 2.7, 1.4, 0.8, "User")
    cases = [
        ("Register", 2.5, 1.2),
        ("Login", 4.9, 1.2),
        ("Browse Products", 7.3, 1.2),
        ("Upload Photo", 2.5, 3.0),
        ("Try Virtual Outfit", 4.9, 3.0),
        ("View Recommendation", 7.3, 3.0),
        ("Add to Cart", 3.7, 4.8),
        ("Place Order", 6.1, 4.8),
    ]
    for txt, x, y in cases:
        o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(2.0), Inches(0.75))
        o.fill.solid()
        o.fill.fore_color.rgb = RGBColor(243, 244, 246)
        o.line.color.rgb = RGBColor(107, 114, 128)
        o.text_frame.text = txt
        o.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def sequence(prs):
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Sequence Diagram"
    cols = [("User", 0.6), ("Frontend", 2.8), ("API", 5.0), ("AI Service", 7.2), ("Database", 9.2)]
    for name, x in cols:
        add_box(s, x, 1.0, 1.6, 0.6, name, 11)
        line = s.shapes.add_shape(MSO_SHAPE.LINE_INVERSE, Inches(x + 0.8), Inches(1.6), Inches(0.01), Inches(4.8))
        line.line.color.rgb = RGBColor(156, 163, 175)

    rows = [
        (0.8, 2.8, "Select item + upload image"),
        (2.8, 5.0, "POST try-on request"),
        (5.0, 7.2, "Forward inference request"),
        (7.2, 5.0, "Return result + size"),
        (5.0, 9.2, "Store metadata"),
        (5.0, 2.8, "Send response"),
        (2.8, 0.8, "Render result"),
    ]
    y = 2.0
    for sx, ex, text in rows:
        add_arrow(s, sx, y, max(ex - sx - 0.1, 0.55), 0.22, text)
        y += 0.58


def activity(prs):
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Activity Diagram"
    steps = [
        "Start",
        "Open app",
        "Choose product",
        "Upload photo",
        "Validate image",
        "Send to API",
        "Run AI inference",
        "Build response",
        "Display try-on output",
        "End",
    ]
    y = 0.9
    for i, step in enumerate(steps):
        add_box(s, 3.4, y, 4.5, 0.42, step, 11)
        if i < len(steps) - 1:
            add_down(s, 5.35, y + 0.42, 0.5, 0.28)
        y += 0.7


def component(prs):
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Component Diagram"
    add_box(s, 0.6, 1.2, 2.5, 0.8, "Client UI")
    add_box(s, 0.6, 2.3, 2.5, 0.8, "HTTP Client")

    add_box(s, 3.8, 1.0, 2.6, 0.8, "Auth Controller")
    add_box(s, 3.8, 2.1, 2.6, 0.8, "Product Controller")
    add_box(s, 3.8, 3.2, 2.6, 0.8, "Inference Controller")
    add_box(s, 3.8, 4.3, 2.6, 0.8, "Cart/Order Controller")

    add_box(s, 7.2, 1.5, 2.0, 0.8, "MongoDB")
    add_box(s, 7.2, 3.4, 2.0, 0.8, "Flask AI API")
    add_box(s, 9.5, 3.4, 2.4, 0.8, "Inference Engine")

    add_arrow(s, 3.1, 2.55, 0.6, 0.25)
    add_arrow(s, 6.5, 1.85, 0.6, 0.25)
    add_arrow(s, 6.5, 3.55, 0.6, 0.25)
    add_arrow(s, 9.2, 3.55, 0.25, 0.25)


def deployment(prs):
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Deployment Diagram"
    add_box(s, 0.8, 2.4, 2.3, 1.0, "User Device\nBrowser")
    add_box(s, 3.6, 1.0, 2.8, 1.0, "Frontend Host\nStatic Web App")
    add_box(s, 3.6, 3.8, 2.8, 1.0, "Backend Host\nNode.js + Express")
    add_box(s, 7.2, 1.0, 2.6, 1.0, "Database Host\nMongoDB")
    add_box(s, 7.2, 3.8, 2.6, 1.0, "AI Host\nFlask Service")

    add_arrow(s, 3.1, 2.7, 0.45, 0.25)
    add_arrow(s, 5.0, 2.1, 0.25, 1.6)
    add_arrow(s, 6.5, 1.35, 0.65, 0.25)
    add_arrow(s, 6.5, 4.15, 0.65, 0.25)


def generate(path="AI_Virtual_TryOn_Diagrams_NoGender.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title(prs, "AI Virtual Try-On", "Required Diagrams (No Gender / No DeepFace)")
    architecture(prs)
    use_case(prs)
    sequence(prs)
    activity(prs)
    component(prs)
    deployment(prs)

    prs.save(path)
    print(f"Generated: {path}")


if __name__ == "__main__":
    generate()
